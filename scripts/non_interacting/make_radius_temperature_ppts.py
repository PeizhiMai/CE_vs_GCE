#!/usr/bin/env python3
"""Create per-radius PowerPoint decks directly with python-pptx.

Policy note:
- Edit PowerPoint files directly in this script.
- Do not use Keynote/AppleScript as an intermediate export path.
"""
from __future__ import annotations

import argparse
import csv
from pathlib import Path

from PIL import Image
from pptx import Presentation
from pptx.enum.text import MSO_ANCHOR, MSO_AUTO_SIZE, PP_ALIGN
from pptx.util import Inches, Pt


ROOT = Path("/Users/cosdis/Desktop/projects/CE_GCE")
DEFAULT_INPUT_DIR = ROOT / "results" / "non_interacting" / "spinless_disk_comparison_by_radius"

SLIDE_WIDTH_IN = 13.333333
SLIDE_HEIGHT_IN = 7.5


def parse_args() -> argparse.Namespace:
    parser = argparse.ArgumentParser(
        description="Create one PowerPoint per radius from all temperature-axis comparison figures."
    )
    parser.add_argument(
        "--input-dir",
        type=Path,
        default=DEFAULT_INPUT_DIR,
        help="Directory containing radius_* folders with *vs_temperature.png files.",
    )
    parser.add_argument(
        "--radii",
        default="auto",
        help="Comma-separated radii list, or 'auto'/'all' to use detected radius_* folders.",
    )
    return parser.parse_args()


def detect_radii(input_dir: Path) -> list[str]:
    radii = []
    for path in sorted(input_dir.glob("radius_*")):
        if not path.is_dir():
            continue
        radius = path.name.removeprefix("radius_")
        try:
            int(radius)
        except ValueError:
            continue
        radii.append(radius)
    return sorted(radii, key=int)


def px_to_in_x(value: float) -> float:
    return value * SLIDE_WIDTH_IN / 1920.0


def px_to_in_y(value: float) -> float:
    return value * SLIDE_HEIGHT_IN / 1080.0


def read_first_row(path: Path) -> list[str]:
    with path.open(newline="") as handle:
        reader = csv.reader(handle, delimiter="\t")
        next(reader)
        return next(reader)


def add_centered_title(slide, text: str) -> None:
    textbox = slide.shapes.add_textbox(
        Inches(px_to_in_x(80)),
        Inches(px_to_in_y(4)),
        Inches(px_to_in_x(1760)),
        Inches(px_to_in_y(56)),
    )
    text_frame = textbox.text_frame
    text_frame.clear()
    text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
    text_frame.auto_size = MSO_AUTO_SIZE.NONE
    paragraph = text_frame.paragraphs[0]
    paragraph.alignment = PP_ALIGN.CENTER
    run = paragraph.add_run()
    run.text = text
    run.font.size = Pt(28)
    run.font.bold = False


def image_height_in(path: Path, width_in: float) -> float:
    with Image.open(path) as img:
        width_px, height_px = img.size
    return width_in * height_px / width_px


def add_picture(slide, path: Path, x_px: float, y_px: float, width_px: float) -> None:
    slide.shapes.add_picture(
        str(path),
        Inches(px_to_in_x(x_px)),
        Inches(px_to_in_y(y_px)),
        width=Inches(px_to_in_x(width_px)),
    )


def make_deck(radius_dir: Path) -> Path:
    radius_dir = radius_dir.resolve()
    images = [
        radius_dir / "energy_vs_temperature.png",
        radius_dir / "specific_heat_vs_temperature.png",
        radius_dir / "compressibility_vs_temperature.png",
        radius_dir / "entropy_vs_temperature.png",
    ]
    missing = [path for path in images if not path.exists()]
    if missing:
        raise FileNotFoundError(f"Missing images for {radius_dir.name}: {missing}")

    first_row = read_first_row(radius_dir / "energy_vs_beta.tsv")
    radius = radius_dir.name.replace("radius_", "")
    nsites = first_row[2]
    canonical_n = first_row[3]
    title = (
        f"Disk, OBC, R={radius}, V={nsites}, N={canonical_n}, "
        f"GCE tuned to <N>={canonical_n}"
    )

    prs = Presentation()
    prs.slide_width = Inches(SLIDE_WIDTH_IN)
    prs.slide_height = Inches(SLIDE_HEIGHT_IN)
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    add_centered_title(slide, title)

    positions = [
        (81.0, 58.0),
        (1021.0, 58.0),
        (81.0, 560.0),
        (1021.0, 560.0),
    ]
    image_width_px = 837.9

    bottom_limit_in = SLIDE_HEIGHT_IN - 0.1
    height_in = image_height_in(images[0], px_to_in_x(image_width_px))
    bottom_in = px_to_in_y(560.0) + height_in
    if bottom_in > bottom_limit_in:
        available_height_in = bottom_limit_in - px_to_in_y(560.0)
        image_width_px *= available_height_in / height_in

    for path, (x_px, y_px) in zip(images, positions, strict=True):
        add_picture(slide, path, x_px, y_px, image_width_px)

    out_path = radius_dir / f"{radius_dir.name}_temperature_figures.pptx"
    prs.save(out_path)
    return out_path


def main() -> None:
    args = parse_args()
    radii_arg = args.radii.strip().lower()
    if radii_arg in {"auto", "all"}:
        radii = detect_radii(args.input_dir)
    else:
        radii = [item.strip() for item in args.radii.split(",") if item.strip()]

    for radius in radii:
        radius_dir = args.input_dir / f"radius_{radius}"
        out_path = make_deck(radius_dir)
        print(f"Created {out_path}")


if __name__ == "__main__":
    main()
