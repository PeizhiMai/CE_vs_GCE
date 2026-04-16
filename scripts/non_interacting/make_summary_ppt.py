#!/usr/bin/env python3
"""Build summary PowerPoint decks directly with python-pptx.

Policy note:
- Edit PowerPoint files directly in this script.
- Do not use Keynote/AppleScript as an intermediate export path.
"""
from __future__ import annotations

import argparse
import csv
from pathlib import Path

from PIL import Image
from pptx import Presentation
from pptx.enum.text import MSO_ANCHOR, MSO_AUTO_SIZE, PP_ALIGN
from pptx.util import Inches, Pt


ROOT = Path("/Users/cosdis/Desktop/projects/CE_GCE")
DEFAULT_INPUT_DIR = ROOT / "results" / "non_interacting" / "spinless_disk_comparison_by_radius"
DEFAULT_OUTPUT = DEFAULT_INPUT_DIR / "summary.pptx"
DEFAULT_SQUARE_OBC_DIR = ROOT / "results" / "non_interacting" / "spinless_square_obc_comparison_by_L"
DEFAULT_SQUARE_PBC_DIR = ROOT / "results" / "non_interacting" / "spinless_square_pbc_comparison_by_L"

SLIDE_WIDTH_IN = 13.333333
SLIDE_HEIGHT_IN = 7.5
REFERENCE_PANEL_HEIGHT_PX = 513.0


def parse_args() -> argparse.Namespace:
    parser = argparse.ArgumentParser(description="Build summary.pptx from selected radius folders.")
    parser.add_argument(
        "--input-dir",
        type=Path,
        default=DEFAULT_INPUT_DIR,
        help="Directory containing radius_* folders.",
    )
    parser.add_argument(
        "--radii",
        default="auto",
        help="Comma-separated radii to use for slides 1..N, or 'auto'/'all' to use detected radius_* folders.",
    )
    parser.add_argument(
        "--output",
        type=Path,
        default=DEFAULT_OUTPUT,
        help="Output PowerPoint path.",
    )
    parser.add_argument(
        "--section-slide",
        action="append",
        default=[],
        help="Insert a section slide as INDEX:TITLE using the manually added section-slide style.",
    )
    parser.add_argument(
        "--square-obc-dir",
        type=Path,
        default=DEFAULT_SQUARE_OBC_DIR,
        help="Directory containing Square OBC L_* folders.",
    )
    parser.add_argument(
        "--square-obc-after",
        type=int,
        default=None,
        help="Insert Square OBC summary slides immediately after this slide index.",
    )
    parser.add_argument(
        "--square-obc-Ls",
        default="",
        help="Comma-separated Square OBC entries to add, e.g. 4,8,16,32,72,71_N_2521.",
    )
    parser.add_argument(
        "--square-pbc-dir",
        type=Path,
        default=DEFAULT_SQUARE_PBC_DIR,
        help="Directory containing Square PBC L_* folders.",
    )
    parser.add_argument(
        "--square-pbc-after",
        type=int,
        default=None,
        help="Insert Square PBC summary slides immediately after this slide index.",
    )
    parser.add_argument(
        "--square-pbc-Ls",
        default="",
        help="Comma-separated Square PBC entries to add, e.g. 4,8,16,32,72,71_N_2521.",
    )
    return parser.parse_args()


def detect_radii(input_dir: Path) -> list[str]:
    radii = []
    for path in sorted(input_dir.glob("radius_*")):
        if not path.is_dir():
            continue
        radius = path.name.removeprefix("radius_")
        try:
            int(radius)
        except ValueError:
            continue
        radii.append(radius)
    return sorted(radii, key=int)


def read_first_row(path: Path) -> list[str]:
    with path.open(newline="") as handle:
        reader = csv.reader(handle, delimiter="\t")
        next(reader)
        return next(reader)


def px_to_in_x(value: float) -> float:
    return value * SLIDE_WIDTH_IN / 1920.0


def px_to_in_y(value: float) -> float:
    return value * SLIDE_HEIGHT_IN / 1080.0


def add_centered_title(slide, text: str) -> None:
    textbox = slide.shapes.add_textbox(
        Inches(px_to_in_x(80)),
        Inches(px_to_in_y(4)),
        Inches(px_to_in_x(1760)),
        Inches(px_to_in_y(56)),
    )
    text_frame = textbox.text_frame
    text_frame.clear()
    text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
    text_frame.auto_size = MSO_AUTO_SIZE.NONE
    paragraph = text_frame.paragraphs[0]
    paragraph.alignment = PP_ALIGN.CENTER
    run = paragraph.add_run()
    run.text = text
    run.font.size = Pt(28)
    run.font.bold = False


def add_section_slide(prs: Presentation, text: str) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    textbox = slide.shapes.add_textbox(
        Inches(3.333),
        Inches(3.189),
        Inches(6.667),
        Inches(0.774),
    )
    text_frame = textbox.text_frame
    text_frame.clear()
    text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
    text_frame.auto_size = MSO_AUTO_SIZE.NONE
    paragraph = text_frame.paragraphs[0]
    paragraph.alignment = PP_ALIGN.CENTER
    run = paragraph.add_run()
    run.text = text
    run.font.size = Pt(40)
    run.font.bold = False


def add_picture(slide, path: Path, x_px: float, y_px: float, height_px: float) -> None:
    slide.shapes.add_picture(
        str(path),
        Inches(px_to_in_x(x_px)),
        Inches(px_to_in_y(y_px)),
        height=Inches(px_to_in_y(height_px)),
    )


def build_slide(prs: Presentation, title: str, images: list[Path]) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_centered_title(slide, title)
    positions = [
        (81.0, 58.0),
        (1021.0, 58.0),
        (81.0, 560.0),
        (1021.0, 560.0),
    ]
    for path, (x_px, y_px) in zip(images, positions, strict=True):
        add_picture(slide, path, x_px, y_px, REFERENCE_PANEL_HEIGHT_PX)


def parse_section_slides(specs: list[str]) -> dict[int, str]:
    section_slides: dict[int, str] = {}
    for spec in specs:
        index_text, sep, title = spec.partition(":")
        if not sep or not title.strip():
            raise ValueError(f"Invalid --section-slide '{spec}'. Expected INDEX:TITLE.")
        index = int(index_text)
        if index < 1:
            raise ValueError(f"Invalid --section-slide index '{index}'. Must be >= 1.")
        if index in section_slides:
            raise ValueError(f"Duplicate --section-slide index '{index}'.")
        section_slides[index] = title.strip()
    return section_slides


def parse_square_entries(spec: str) -> list[str]:
    return [item.strip() for item in spec.split(",") if item.strip()]


def load_square_payload(input_dir: Path, entry: str, boundary_label: str) -> tuple[str, list[Path]]:
    folder = input_dir / f"L_{entry}"
    if not folder.exists():
        raise FileNotFoundError(f"Missing Square {boundary_label} folder for entry '{entry}': {folder}")

    images = [
        folder / "energy_vs_temperature.png",
        folder / "specific_heat_vs_temperature.png",
        folder / "compressibility_vs_temperature.png",
        folder / "entropy_vs_temperature.png",
    ]
    missing = [path for path in images if not path.exists()]
    if missing:
        raise FileNotFoundError(f"Missing Square {boundary_label} images for entry '{entry}': {missing}")

    metadata = {}
    with (folder / "metadata.toml").open() as handle:
        for raw_line in handle:
            line = raw_line.strip()
            if " = " not in line:
                continue
            key, value = line.split(" = ", 1)
            metadata[key] = value.strip().strip('"')

    lattice_size = metadata["L"]
    nsites = metadata["nsites"]
    canonical_n = metadata["canonical_n"]
    title = (
        f"Square, {boundary_label}, L={lattice_size}, V={nsites}, N={canonical_n}, "
        f"GCE tuned to <N>={canonical_n}"
    )
    return title, images


def main() -> None:
    args = parse_args()
    radii_arg = args.radii.strip().lower()
    if radii_arg in {"auto", "all"}:
        radii = detect_radii(args.input_dir)
    else:
        radii = [item.strip() for item in args.radii.split(",") if item.strip()]
    section_slides = parse_section_slides(args.section_slide)
    square_obc_entries = parse_square_entries(args.square_obc_Ls)
    square_pbc_entries = parse_square_entries(args.square_pbc_Ls)

    prs = Presentation()
    prs.slide_width = Inches(SLIDE_WIDTH_IN)
    prs.slide_height = Inches(SLIDE_HEIGHT_IN)
    if prs.slides:
        del prs.slides._sldIdLst[0]

    slide_payloads = []
    for radius in radii:
        radius_dir = (args.input_dir / f"radius_{radius}").resolve()
        images = [
            radius_dir / "energy_vs_temperature.png",
            radius_dir / "specific_heat_vs_temperature.png",
            radius_dir / "compressibility_vs_temperature.png",
            radius_dir / "entropy_vs_temperature.png",
        ]
        missing = [path for path in images if not path.exists()]
        if missing:
            raise FileNotFoundError(f"Missing images for radius {radius}: {missing}")
        first_row = read_first_row(radius_dir / "energy_vs_beta.tsv")
        nsites = first_row[2]
        canonical_n = first_row[3]
        title = (
            f"Disk, OBC, R={radius}, V={nsites}, N={canonical_n}, "
            f"GCE tuned to <N>={canonical_n}"
        )
        slide_payloads.append((title, images))

    square_obc_payloads = []
    if square_obc_entries:
        if args.square_obc_after is None:
            raise ValueError("--square-obc-after is required when --square-obc-Ls is provided.")
        square_obc_payloads = [load_square_payload(args.square_obc_dir, entry, "OBC") for entry in square_obc_entries]

    square_pbc_payloads = []
    if square_pbc_entries:
        if args.square_pbc_after is None:
            raise ValueError("--square-pbc-after is required when --square-pbc-Ls is provided.")
        square_pbc_payloads = [load_square_payload(args.square_pbc_dir, entry, "PBC") for entry in square_pbc_entries]

    final_items: list[tuple[str, object]] = [("content", payload) for payload in slide_payloads]
    for index, title in sorted(section_slides.items()):
        final_items.insert(index - 1, ("section", title))

    if square_obc_payloads:
        insert_at = args.square_obc_after
        for offset, payload in enumerate(square_obc_payloads):
            final_items.insert(insert_at + offset, ("content", payload))

    if square_pbc_payloads:
        insert_at = args.square_pbc_after
        for offset, payload in enumerate(square_pbc_payloads):
            final_items.insert(insert_at + offset, ("content", payload))

    for item_type, payload in final_items:
        if item_type == "section":
            add_section_slide(prs, payload)
            continue
        title, images = payload
        build_slide(prs, title, images)

    prs.save(args.output)
    print(f"Created {args.output}")


if __name__ == "__main__":
    main()
